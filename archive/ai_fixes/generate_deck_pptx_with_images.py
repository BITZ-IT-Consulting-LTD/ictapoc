import os
import re
import subprocess
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def extract_mermaid(content, pattern="AS-IS"):
    """Extract mermaid code block based on heading"""
    blocks = re.split(r'```mermaid\n', content)
    for i, block in enumerate(blocks):
        if i == 0: continue
        mermaid_code = block.split('```')[0]
        # Very rough heuristic: if we want AS-IS, we assume it's the first diagram, TO-BE is the second
        return mermaid_code
    return None

def extract_all_mermaids(content):
    matches = re.findall(r'```mermaid\n(.*?)```', content, re.DOTALL)
    return matches

def main():
    mda_dir = "/Users/mac/ictapoc/mdas/docs_final/priority_mdas/"
    md_files = [f for f in os.listdir(mda_dir) if f.endswith(".md") and f != "Priority_MDAs_Justification_Matrix.md" and f != "temp_master_combined.md"]
    
    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "KDEAP [16] - Proof of Concept v2"
    subtitle.text = "Digital Service Architecture Programme\nTransforming Government Services"
    
    # Create temp dir for mermaid images
    img_dir = os.path.join(mda_dir, "mermaid_exports")
    os.makedirs(img_dir, exist_ok=True)
    
    for filename in sorted(md_files):
        filepath = os.path.join(mda_dir, filename)
        with open(filepath, 'r') as f:
            content = f.read()
            
        # Extract MDA Name
        mda_match = re.search(r'- \*\*Ministry/Department/Agency \(MDA\):\*\* (.*)', content)
        department_match = re.search(r'- \*\*Department:\*\* (.*)', content)
        authority_match = re.search(r'- \*\*Authority:\*\* (.*)', content)
        office_match = re.search(r'- \*\*Office:\*\* (.*)', content)
        
        mda_name = ""
        if department_match: mda_name = department_match.group(1)
        elif authority_match: mda_name = authority_match.group(1)
        elif office_match: mda_name = office_match.group(1)
        elif mda_match: mda_name = mda_match.group(1)
        else: mda_name = filename.replace(".md", "").replace("_", " ")
        
        # Extract Process Name
        process_match = re.search(r'- \*\*Process Name:\*\* (.*)', content)
        process_name = process_match.group(1) if process_match else "Core Service Delivery"
        
        # Extract Pain Points
        pain_points_section = re.search(r'### Pain Points\n(.*?)(?=\n### |\n## )', content, re.DOTALL)
        pain_points = []
        if pain_points_section:
             items = re.findall(r'- (.*)', pain_points_section.group(1))
             pain_points = [item for item in items[:3]]
        
        # Extract Opportunities/TO-BE Narrative
        opps_section = re.search(r'### Opportunities\n(.*?)(?=\n### |\n## )', content, re.DOTALL)
        opps = []
        if opps_section:
            items = re.findall(r'- (.*)', opps_section.group(1))
            opps = [item for item in items[:3]]

        # Generate Slide 1: Text Summary
        slide_layout = prs.slide_layouts[6] # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        txbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        tf = txbox.text_frame
        p = tf.add_paragraph()
        p.text = mda_name
        p.font.bold = True
        p.font.size = Pt(28)
        p.font.color.rgb = RGBColor(0, 51, 102)
        
        p2 = tf.add_paragraph()
        p2.text = process_name
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(100, 100, 100)
        
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4.2), Inches(4.5))
        tf_left = left_box.text_frame
        tf_left.word_wrap = True
        p_title = tf_left.add_paragraph()
        p_title.text = "Current Pain Points (AS-IS)"
        p_title.font.bold = True
        p_title.font.size = Pt(16)
        p_title.font.color.rgb = RGBColor(200, 0, 0)
        for point in pain_points:
            p = tf_left.add_paragraph()
            p.text = "• " + point
            p.font.size = Pt(12)
            p.space_before = Pt(6)
            
        right_box = slide.shapes.add_textbox(Inches(5.3), Inches(2), Inches(4.2), Inches(4.5))
        tf_right = right_box.text_frame
        tf_right.word_wrap = True
        p_title2 = tf_right.add_paragraph()
        p_title2.text = "Digital Transformation (TO-BE)"
        p_title2.font.bold = True
        p_title2.font.size = Pt(16)
        p_title2.font.color.rgb = RGBColor(0, 150, 0)
        for opp in opps:
            p = tf_right.add_paragraph()
            p.text = "• " + opp
            p.font.size = Pt(12)
            p.space_before = Pt(6)

        # Extract mermaids
        mermaids = extract_all_mermaids(content)
        
        for idx, mermaid_code in enumerate(mermaids[:2]): # Max 2 (AS-IS, TO-BE)
            stage = "AS-IS Architecture" if idx == 0 else "TO-BE Architecture"
            
            # Save mermaid code to temp file
            mmd_path = os.path.join(img_dir, f"temp_{idx}.mmd")
            png_path = os.path.join(img_dir, f"temp_{idx}.png")
            
            with open(mmd_path, 'w') as mf:
                mf.write(mermaid_code)
                
            # Render to PNG using mmdc
            try:
                subprocess.run(["npx", "-p", "@mermaid-js/mermaid-cli", "mmdc", "-i", mmd_path, "-o", png_path, "-b", "transparent", "-t", "default"], check=True, capture_output=True)
                
                # Add slide for diagram
                slide = prs.slides.add_slide(slide_layout)
                
                # Title
                txbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
                tf = txbox.text_frame
                p = tf.add_paragraph()
                p.text = f"{mda_name} - {stage}"
                p.font.bold = True
                p.font.size = Pt(24)
                p.font.color.rgb = RGBColor(0, 51, 102)
                
                # Add image
                # Try to fit it nicely
                try:
                    slide.shapes.add_picture(png_path, Inches(0.5), Inches(1.2), width=Inches(9))
                except Exception as e:
                    print(f"Could not add image {png_path}: {e}")
                    
            except Exception as e:
                print(f"Failed to generate mermaid for {filename}, idx {idx}: {e}")

    out_path = '/Users/mac/ictapoc/mdas/docs_final/priority_mdas/KDEAP_PoC_v2_Comprehensive.pptx'
    prs.save(out_path)
    print(f"Comprehensive PowerPoint generated at {out_path}")

if __name__ == '__main__':
    main()
