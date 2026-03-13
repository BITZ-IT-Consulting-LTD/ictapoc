import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re

def create_slide_deck():
    prs = Presentation()
    
    # Customizing slide master (optional, simplified here)
    # Using layout 0 (Title) for the main title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "KDEAP [16] - Proof of Concept v2"
    subtitle.text = "Digital Service Architecture Programme\nTransforming Government Services"
    
    # Read the markdown file we generated earlier
    with open('/Users/mac/ictapoc/mdas/docs_final/priority_mdas/KDEAP_PoC_v2_Deck.md', 'r') as f:
        content = f.read()
        
    # Split by MDA entries
    sections = re.split(r'\n---\n\n', content)
    
    for section in sections:
        if not section.strip() or section.startswith('# KDEAP'):
            continue
            
        lines = section.strip().split('\n')
        
        mda_name = ""
        process_name = ""
        pain_points = []
        opportunities = []
        
        current_list = None
        
        for line in lines:
            if line.startswith('## '):
                mda_name = line.replace('## ', '').strip()
            elif line.startswith('### '):
                process_name = line.replace('### ', '').strip()
            elif line.startswith('**Current Pain Points:**'):
                current_list = pain_points
            elif line.startswith('**Digital Transformation (TO-BE):**'):
                current_list = opportunities
            elif line.startswith('  * '):
                if current_list is not None:
                    current_list.append(line.replace('  * ', '').strip())
        
        # Create slide for this MDA (using layout 1 - Title and Content)
        # But we'll use a blank layout to customize the boxes
        if not mda_name: continue
        
        slide_layout = prs.slide_layouts[6] # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        txbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        tf = txbox.text_frame
        p = tf.add_paragraph()
        p.text = mda_name
        p.font.bold = True
        p.font.size = Pt(32)
        p.font.color.rgb = RGBColor(0, 51, 102)
        
        # Subtitle (Process Name)
        p2 = tf.add_paragraph()
        p2.text = process_name
        p2.font.size = Pt(20)
        p2.font.color.rgb = RGBColor(100, 100, 100)
        
        # Left Box (Pain Points)
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4.2), Inches(4.5))
        tf_left = left_box.text_frame
        tf_left.word_wrap = True
        p_title = tf_left.add_paragraph()
        p_title.text = "Current Pain Points (AS-IS)"
        p_title.font.bold = True
        p_title.font.size = Pt(18)
        p_title.font.color.rgb = RGBColor(200, 0, 0)
        
        for point in pain_points:
            p = tf_left.add_paragraph()
            p.text = "• " + point
            p.font.size = Pt(14)
            p.level = 0
            p.space_before = Pt(6)
            
        # Right Box (Opportunities)
        right_box = slide.shapes.add_textbox(Inches(5.3), Inches(2), Inches(4.2), Inches(4.5))
        tf_right = right_box.text_frame
        tf_right.word_wrap = True
        p_title2 = tf_right.add_paragraph()
        p_title2.text = "Digital Transformation (TO-BE)"
        p_title2.font.bold = True
        p_title2.font.size = Pt(18)
        p_title2.font.color.rgb = RGBColor(0, 150, 0)
        
        for opp in opportunities:
            p = tf_right.add_paragraph()
            p.text = "• " + opp
            p.font.size = Pt(14)
            p.level = 0
            p.space_before = Pt(6)

    prs.save('/Users/mac/ictapoc/mdas/docs_final/priority_mdas/KDEAP_PoC_v2_Deck.pptx')
    print("PowerPoint generated at /Users/mac/ictapoc/mdas/docs_final/priority_mdas/KDEAP_PoC_v2_Deck.pptx")

if __name__ == '__main__':
    create_slide_deck()
